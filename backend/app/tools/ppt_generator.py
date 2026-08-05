import io
import logging
from typing import Dict, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

logger = logging.getLogger("synovia.tools.ppt_generator")

def clean_str(val: Any) -> str:
    if val is None:
        return ""
    if isinstance(val, (int, float, bool)):
        return str(val)
    if isinstance(val, dict):
        parts = [f"{k}: {clean_str(v)}" for k, v in val.items() if v]
        return " • ".join(parts)
    if isinstance(val, list):
        return ", ".join(clean_str(v) for v in val if v)
    return str(val).replace("■", "").strip()

def parse_score(val: Any, default: int = 80) -> int:
    if isinstance(val, (int, float)):
        return int(val)
    if isinstance(val, dict):
        return parse_score(val.get("score") or val.get("value"), default)
    if isinstance(val, str):
        digits = "".join(filter(str.isdigit, val))
        return int(digits) if digits else default
    return default

class PPTReportGenerator:
    """
    Generates an executive 16:9 PowerPoint (.pptx) Pitch Deck from blueprint data.
    """
    def __init__(self):
        # Premium Dark Palette (Slate 950, Indigo 600, Emerald 400)
        self.bg_color = RGBColor(15, 23, 42)        # Slate 900
        self.card_bg = RGBColor(30, 41, 59)        # Slate 800
        self.text_white = RGBColor(255, 255, 255)
        self.text_indigo = RGBColor(129, 140, 248)  # Indigo 400
        self.text_emerald = RGBColor(52, 211, 153)  # Emerald 400
        self.text_slate = RGBColor(148, 163, 184)   # Slate 400

    def create_deck(self, blueprint: Dict[str, Any]) -> bytes:
        prs = Presentation()
        # Set 16:9 widescreen layout
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        idea = clean_str(blueprint.get("idea", "Startup Blueprint")).title()
        research = blueprint.get("research", {})
        competitor = blueprint.get("competitor", {})
        product = blueprint.get("product", {})
        validation = blueprint.get("validation", {})
        roadmap = blueprint.get("roadmap", {})
        pitch = blueprint.get("pitch", {})

        # Slide 1: Cover Slide
        self._add_cover_slide(prs, idea, blueprint.get("executive_summary", ""))

        # Slide 2: Market Opportunity & TAM
        self._add_market_slide(prs, idea, research)

        # Slide 3: Competitors & Market Gaps
        self._add_competitor_slide(prs, idea, competitor)

        # Slide 4: MVP Product Specification
        self._add_product_slide(prs, idea, product)

        # Slide 5: Validation & Strategy Report (VC Scores & Verdict)
        self._add_validation_slide(prs, idea, validation)

        # Slide 6: 4-Week Execution Roadmap
        self._add_roadmap_slide(prs, idea, roadmap)

        # Slide 7: Investor Pitch & Business Model
        self._add_pitch_slide(prs, idea, pitch)

        buffer = io.BytesIO()
        prs.save(buffer)
        return buffer.getvalue()

    def _set_slide_background(self, slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.bg_color

    def _add_header(self, slide, title_text: str, subtitle_text: str = ""):
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.text_white

        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.size = Pt(13)
            p2.font.color.rgb = self.text_indigo

    def _add_cover_slide(self, prs, idea: str, summary: str):
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        self._set_slide_background(slide)

        # Card shape background
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.2), Inches(11.333), Inches(5.1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.card_bg
        shape.line.color.rgb = self.text_indigo

        txBox = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10.333), Inches(4.0))
        tf = txBox.text_frame
        tf.word_wrap = True

        p0 = tf.paragraphs[0]
        p0.text = "SYNOVIA AI — INVESTOR PITCH DECK"
        p0.font.size = Pt(14)
        p0.font.bold = True
        p0.font.color.rgb = self.text_emerald

        p1 = tf.add_paragraph()
        p1.text = idea
        p1.font.size = Pt(36)
        p1.font.bold = True
        p1.font.color.rgb = self.text_white
        p1.space_before = Pt(10)

        p2 = tf.add_paragraph()
        p2.text = summary or "Autonomous AI Multi-Agent Generated Startup Blueprint & Strategy."
        p2.font.size = Pt(14)
        p2.font.color.rgb = self.text_slate
        p2.space_before = Pt(15)

    def _add_market_slide(self, prs, idea: str, research: Dict[str, Any]):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_slide_background(slide)
        self._add_header(slide, "Market Size & Opportunity (TAM/SAM/SOM)", f"Project: {idea}")

        tam = clean_str(research.get("market_size", {}).get("tam", "N/A"))
        sam = clean_str(research.get("market_size", {}).get("sam", "N/A"))
        som = clean_str(research.get("market_size", {}).get("som", "N/A"))

        # 3 Market Box Shapes
        boxes = [("TAM (Total Market)", tam, Inches(0.8)), ("SAM (Serviceable Market)", sam, Inches(4.8)), ("SOM (Target Reachable)", som, Inches(8.8))]
        for label, val, left in boxes:
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.8), Inches(3.7), Inches(1.8))
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.card_bg
            shape.line.color.rgb = self.text_indigo

            tf = shape.text_frame
            tf.word_wrap = True
            p1 = tf.paragraphs[0]
            p1.text = label
            p1.font.size = Pt(12)
            p1.font.bold = True
            p1.font.color.rgb = self.text_indigo

            p2 = tf.add_paragraph()
            p2.text = val
            p2.font.size = Pt(14)
            p2.font.color.rgb = self.text_white
            p2.space_before = Pt(8)

        # Customer Pain Points Card
        shape2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.9), Inches(11.7), Inches(2.8))
        shape2.fill.solid()
        shape2.fill.fore_color.rgb = self.card_bg
        shape2.line.color.rgb = self.text_slate

        tf2 = shape2.text_frame
        tf2.word_wrap = True
        p_title = tf2.paragraphs[0]
        p_title.text = "Customer Pain Points & Target Needs"
        p_title.font.size = Pt(14)
        p_title.font.bold = True
        p_title.font.color.rgb = self.text_emerald

        for pain in research.get("customer_pain_points", [])[:4]:
            p = tf2.add_paragraph()
            p.text = f"•  {clean_str(pain)}"
            p.font.size = Pt(12)
            p.font.color.rgb = self.text_white
            p.space_before = Pt(4)

    def _add_competitor_slide(self, prs, idea: str, competitor: Dict[str, Any]):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_slide_background(slide)
        self._add_header(slide, "Competitor Intelligence & Market Moat", f"Project: {idea}")

        left = Inches(0.8)
        for comp in (competitor.get("competitors", [])[:2]):
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.8), Inches(5.7), Inches(4.9))
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.card_bg
            shape.line.color.rgb = self.text_indigo

            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"Competitor: {clean_str(comp.get('name', 'N/A'))}"
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.text_emerald

            p_s = tf.add_paragraph()
            p_s.text = "Strengths:"
            p_s.font.bold = True
            p_s.font.size = Pt(12)
            p_s.font.color.rgb = self.text_indigo
            p_s.space_before = Pt(6)

            for s in comp.get("strengths", [])[:2]:
                p_item = tf.add_paragraph()
                p_item.text = f"• {clean_str(s)}"
                p_item.font.size = Pt(11)
                p_item.font.color.rgb = self.text_white

            p_w = tf.add_paragraph()
            p_w.text = "Weaknesses & Market Gaps:"
            p_w.font.bold = True
            p_w.font.size = Pt(12)
            p_w.font.color.rgb = self.text_slate
            p_w.space_before = Pt(6)

            for w in comp.get("weaknesses", [])[:2]:
                p_item = tf.add_paragraph()
                p_item.text = f"• {clean_str(w)}"
                p_item.font.size = Pt(11)
                p_item.font.color.rgb = self.text_white

            left += Inches(6.0)

    def _add_product_slide(self, prs, idea: str, product: Dict[str, Any]):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_slide_background(slide)
        self._add_header(slide, "MVP Feature Specification", f"Project: {idea}")

        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.card_bg
        shape.line.color.rgb = self.text_indigo

        tf = shape.text_frame
        tf.word_wrap = True

        p0 = tf.paragraphs[0]
        p0.text = "Core MVP Features & User Journey"
        p0.font.size = Pt(16)
        p0.font.bold = True
        p0.font.color.rgb = self.text_emerald

        for feat in product.get("mvp_features", [])[:4]:
            p = tf.add_paragraph()
            feat_name = clean_str(feat.get("name") or feat.get("title") or feat) if isinstance(feat, dict) else clean_str(feat)
            feat_desc = clean_str(feat.get("description") or feat.get("desc") or "") if isinstance(feat, dict) else ""
            p.text = f"•  {feat_name}: {feat_desc}" if feat_desc else f"•  {feat_name}"
            p.font.size = Pt(12)
            p.font.color.rgb = self.text_white
            p.space_before = Pt(8)

    def _add_validation_slide(self, prs, idea: str, validation: Dict[str, Any]):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_slide_background(slide)
        self._add_header(slide, "Validation & Strategy Report (VC Scores)", f"Project: {idea}")

        # 5 Scores Box
        scores = [
            ("Viability", parse_score(validation.get("viability_score"), 82)),
            ("Innovation", parse_score(validation.get("innovation_score"), 78)),
            ("Market Opp.", parse_score(validation.get("market_opportunity_score"), 88)),
            ("Feasibility", parse_score(validation.get("feasibility_score"), 75)),
            ("Scalability", parse_score(validation.get("scalability_score"), 84)),
        ]

        left = Inches(0.8)
        for label, val in scores:
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.8), Inches(2.2), Inches(1.5))
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.card_bg
            shape.line.color.rgb = self.text_emerald

            tf = shape.text_frame
            p1 = tf.paragraphs[0]
            p1.text = label
            p1.font.size = Pt(11)
            p1.font.color.rgb = self.text_slate

            p2 = tf.add_paragraph()
            p2.text = f"{val}/100"
            p2.font.size = Pt(20)
            p2.font.bold = True
            p2.font.color.rgb = self.text_emerald
            p2.space_before = Pt(4)

            left += Inches(2.4)

        # Verdict Card
        shape2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.6), Inches(11.7), Inches(3.2))
        shape2.fill.solid()
        shape2.fill.fore_color.rgb = self.card_bg
        shape2.line.color.rgb = self.text_indigo

        tf2 = shape2.text_frame
        tf2.word_wrap = True

        p_vtitle = tf2.paragraphs[0]
        p_vtitle.text = "VC Mentor Verdict & Recommendations"
        p_vtitle.font.size = Pt(14)
        p_vtitle.font.bold = True
        p_vtitle.font.color.rgb = self.text_white

        p_v = tf2.add_paragraph()
        p_v.text = f"Final Verdict: {clean_str(validation.get('final_verdict', 'STRONG PURSUE'))}"
        p_v.font.size = Pt(13)
        p_v.font.bold = True
        p_v.font.color.rgb = self.text_emerald
        p_v.space_before = Pt(6)

        for rec in validation.get("validation_recommendations", [])[:2]:
            p = tf2.add_paragraph()
            p.text = f"• Recommendation: {clean_str(rec)}"
            p.font.size = Pt(11)
            p.font.color.rgb = self.text_slate
            p.space_before = Pt(4)

    def _add_roadmap_slide(self, prs, idea: str, roadmap: Dict[str, Any]):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_slide_background(slide)
        self._add_header(slide, "4-Week Agile Execution Roadmap", f"Project: {idea}")

        left = Inches(0.8)
        for wk in roadmap.get("schedule", [])[:4]:
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.8), Inches(2.7), Inches(4.9))
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.card_bg
            shape.line.color.rgb = self.text_indigo

            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"Week {clean_str(wk.get('week', 1))}: {clean_str(wk.get('title', ''))}"
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = self.text_emerald

            p_g = tf.add_paragraph()
            p_g.text = f"Focus: {clean_str(wk.get('goals', ''))}"
            p_g.font.size = Pt(10)
            p_g.font.italic = True
            p_g.font.color.rgb = self.text_slate
            p_g.space_before = Pt(4)

            p_dhead = tf.add_paragraph()
            p_dhead.text = "Deliverables:"
            p_dhead.font.bold = True
            p_dhead.font.size = Pt(10)
            p_dhead.font.color.rgb = self.text_indigo
            p_dhead.space_before = Pt(6)

            for d in wk.get("deliverables", [])[:3]:
                p_item = tf.add_paragraph()
                p_item.text = f"• {clean_str(d)}"
                p_item.font.size = Pt(10)
                p_item.font.color.rgb = self.text_white

            left += Inches(3.0)

    def _add_pitch_slide(self, prs, idea: str, pitch: Dict[str, Any]):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_slide_background(slide)
        self._add_header(slide, "60-Second Investor Elevator Pitch", f"Project: {idea}")

        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.card_bg
        shape.line.color.rgb = self.text_emerald

        tf = shape.text_frame
        tf.word_wrap = True

        p0 = tf.paragraphs[0]
        p0.text = "Elevator Pitch Script"
        p0.font.size = Pt(16)
        p0.font.bold = True
        p0.font.color.rgb = self.text_emerald

        p_script = tf.add_paragraph()
        p_script.text = f"\"{clean_str(pitch.get('hackathon_pitch', 'Innovative solution addressing key market pain points with scalable business economics.'))}\""
        p_script.font.size = Pt(14)
        p_script.font.italic = True
        p_script.font.color.rgb = self.text_white
        p_script.space_before = Pt(10)

        p_usphead = tf.add_paragraph()
        p_usphead.text = "Unfair Advantage / USP:"
        p_usphead.font.bold = True
        p_usphead.font.size = Pt(12)
        p_usphead.font.color.rgb = self.text_indigo
        p_usphead.space_before = Pt(15)

        p_usp = tf.add_paragraph()
        p_usp.text = clean_str(pitch.get("usp", "Proprietary design innovation and rapid market execution."))
        p_usp.font.size = Pt(12)
        p_usp.font.color.rgb = self.text_slate

ppt_report_generator = PPTReportGenerator()
